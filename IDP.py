# ==============================
# Tech Doc PROCESSOR
#
# ==============================

import re
import time
import zipfile
import tempfile
import hashlib
from io import BytesIO
from pathlib import Path
from copy import deepcopy
import textwrap
import json

import pandas as pd
import streamlit as st

from docx import Document as DocxDocument
from pptx import Presentation

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from langchain_community.document_loaders import TextLoader, PyPDFLoader

from reportlab.lib.pagesizes import A4
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_LEFT, TA_CENTER
from reportlab.lib import colors

from workflow import build_graph
from core import (
    validate_document_data,
    build_confidence_map,
    classify_exception,
    extract_text_from_pdf_with_ocr_fallback,
    ocr_image_bytes_with_vlm,
)

# ------------------------------
# PAGE CONFIG
# ------------------------------
st.set_page_config("IDP - Professional", layout="wide")
USERS = st.secrets.get("users", {})
MAX_BATCH_FILES = 15

# ------------------------------
# CACHED MODELS
# ------------------------------
@st.cache_resource
def get_llm(api_key, model):
    return ChatOpenAI(model=model, temperature=0, api_key=api_key)


@st.cache_resource
def get_embeddings(api_key):
    return OpenAIEmbeddings(api_key=api_key)

# ------------------------------
# AUTH
# ------------------------------
def validate_api_key(api_key):
    try:
        llm = ChatOpenAI(
            model="gpt-4o-mini",
            temperature=0,
            api_key=api_key
        )
        llm.invoke("Reply with OK")
        return True
    except Exception:
        return False


def login():
    logo_path = Path(__file__).parent / "TDReader.png"
    col1, col2, col3 = st.columns([1, 2, 1])

    with col2:
        if logo_path.exists():
            st.image(logo_path, use_container_width=True)

        st.markdown("### Sign In")
        username = st.text_input("Username")
        password = st.text_input("Password", type="password")
        api_key = st.text_input("OpenAI API Key", type="password")

        if st.button("Login", use_container_width=True):
            if username not in USERS or USERS[username]["password"] != password:
                st.error("Invalid username or password")
                return

            if not api_key:
                st.error("Please enter your OpenAI API key")
                return

            with st.spinner("Validating API key..."):
                if not validate_api_key(api_key):
                    st.error("Invalid OpenAI API key")
                    return

            st.session_state["logged_in"] = True
            st.session_state["user"] = username
            st.session_state["role"] = USERS[username].get("role", "user")
            st.session_state["api_key"] = api_key
            st.rerun()

# ------------------------------
# SESSION INIT
# ------------------------------
DEFAULT_KEYS = {
    "logged_in": False,
    "user": None,
    "role": None,
    "api_key": None,
    "model_choice": "gpt-4o-mini",
    "metrics": {
        "tokens": 0,
        "input_tokens": 0,
        "output_tokens": 0,
        "cost": 0.0,
        "response_times": [],
        "calls": 0
    },
    "doc_costs": {},
    "batch_results": [],
    "exception_queue": [],
    "active_batch_index": 0,
    "batch_processed": False,
    "last_batch_signature": None,
    "show_reprocess_confirm": False,
    "pending_batch_signature": None,
    "batch_total_files": 0,
    "batch_processed_files": 0,
    "batch_current_file": None,
    "batch_file_statuses": [],
    "batch_started_at": None,
    "batch_completed_at": None,
    "batch_elapsed_seconds": 0.0,
    "review_data": None,
    "confidence_map": None,
    "validation_result": None,
    "vectorstore": None,
    "chat_history": [],
    "suggested_questions": [],
    "current_file": None,
    "doc_type": None,
    "full_text": None,
    "auto_result": None,
    "agent_events": [],
    "agent_logs": [],
    "current_step": "Waiting",
    "progress_value": 0,
    "live_step_placeholder": None,
    "live_progress_placeholder": None,
    "live_event_placeholder": None,
    "uploader_key": 0,
    "template_library": [],
    "active_template_index": None,
    "agent_timings": {},
    "active_agent": None,
}

for key, value in DEFAULT_KEYS.items():
    if key not in st.session_state:
        st.session_state[key] = value

if not st.session_state["logged_in"]:
    login()
    st.stop()

# ------------------------------
# HELPERS
# ------------------------------


def reset_run_state():
    st.session_state["review_data"] = None
    st.session_state["confidence_map"] = None
    st.session_state["validation_result"] = None
    st.session_state["vectorstore"] = None
    st.session_state["chat_history"] = []
    st.session_state["suggested_questions"] = []
    st.session_state["current_file"] = None
    st.session_state["doc_type"] = None
    st.session_state["full_text"] = None
    st.session_state["auto_result"] = None
    st.session_state["agent_events"] = []
    st.session_state["agent_logs"] = []
    st.session_state["current_step"] = "Waiting"
    st.session_state["progress_value"] = 0
    st.session_state["live_step_placeholder"] = None
    st.session_state["live_progress_placeholder"] = None
    st.session_state["live_event_placeholder"] = None
    st.session_state["agent_timings"] = {}
    st.session_state["active_agent"] = None

def reset_single_file_state():
    st.session_state["review_data"] = None
    st.session_state["confidence_map"] = None
    st.session_state["validation_result"] = None
    st.session_state["vectorstore"] = None
    st.session_state["chat_history"] = []
    st.session_state["suggested_questions"] = []
    st.session_state["current_file"] = None
    st.session_state["doc_type"] = None
    st.session_state["full_text"] = None
    st.session_state["auto_result"] = None
    st.session_state["agent_events"] = []
    st.session_state["agent_logs"] = []
    st.session_state["current_step"] = "Waiting"
    st.session_state["progress_value"] = 0
    st.session_state["agent_timings"] = {}
    st.session_state["active_agent"] = None

def save_temp_file(uploaded_file):
    suffix = Path(uploaded_file.name).suffix
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(uploaded_file.getvalue())
        return tmp.name


def load_default_resume_template_bytes():
    possible_paths = [
        Path("templates/resume_template.docx"),
        Path("templates:resume_template.docx"),
        Path(__file__).parent / "templates" / "resume_template.docx",
        Path(__file__).parent / "templates:resume_template.docx",
    ]
    for path in possible_paths:
        if path.exists():
            with open(path, "rb") as file:
                return file.read()
    return None


def get_active_template_bytes():
    library = st.session_state.get("template_library", [])
    active_index = st.session_state.get("active_template_index")

    if active_index is not None and 0 <= active_index < len(library):
        return library[active_index]["content"]

    return None


def add_template_to_library(uploaded_template):
    if not uploaded_template:
        return

    content = uploaded_template.getvalue()

    entry = {
        "name": uploaded_template.name,
        "content": content,
    }

    st.session_state.template_library.append(entry)
    st.session_state.active_template_index = len(st.session_state.template_library) - 1


def extract_docx_text(file_path):
    doc = DocxDocument(file_path)
    parts = []

    for paragraph in doc.paragraphs:
        if paragraph.text and paragraph.text.strip():
            parts.append(paragraph.text.strip())

    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    return "\n".join(parts).strip()


def process_file_with_fallback(uploaded_file):
    suffix = Path(uploaded_file.name).suffix.lower()
    uploaded_file.seek(0)

    if suffix in [".png", ".jpg", ".jpeg"]:
        image_bytes = uploaded_file.getvalue()
        mime_type = "image/jpeg" if suffix in [".jpg", ".jpeg"] else "image/png"
        text = ocr_image_bytes_with_vlm(image_bytes, mime_type=mime_type)
        return {
            "documents": [Document(page_content=text)] if text else [],
            "text": text,
            "ocr_used": True,
            "extraction_mode": "image_vlm_ocr",
            "exception_reason": None if text else "OCR failed on image",
        }

    file_path = save_temp_file(uploaded_file)

    try:
        if suffix == ".txt":
            try:
                docs = TextLoader(file_path, encoding="utf-8").load()
            except Exception:
                docs = TextLoader(file_path, encoding="cp1252").load()

            text = "\n".join([d.page_content for d in docs]).strip()
            return {
                "documents": docs,
                "text": text,
                "ocr_used": False,
                "extraction_mode": "plain_text",
                "exception_reason": None,
            }

        if suffix == ".pdf":
            pdf_result = extract_text_from_pdf_with_ocr_fallback(file_path)
            docs = [Document(page_content=pdf_result["text"])] if pdf_result["text"] else []
            return {
                "documents": docs,
                "text": pdf_result["text"],
                "ocr_used": pdf_result["ocr_used"],
                "extraction_mode": pdf_result["extraction_mode"],
                "exception_reason": pdf_result["exception_reason"],
            }

        if suffix == ".docx":
            text = extract_docx_text(file_path)
            docs = [Document(page_content=text)] if text else []
            return {
                "documents": docs,
                "text": text,
                "ocr_used": False,
                "extraction_mode": "docx_text",
                "exception_reason": None if text else "No extractable text in DOCX",
            }

        if suffix == ".pptx":
            prs = Presentation(file_path)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text and shape.text.strip():
                        text_parts.append(shape.text.strip())
            text = "\n".join(text_parts).strip()
            docs = [Document(page_content=text)] if text else []
            return {
                "documents": docs,
                "text": text,
                "ocr_used": False,
                "extraction_mode": "pptx_text",
                "exception_reason": None if text else "No extractable text in PPTX",
            }

        if suffix == ".xlsx":
            excel_file = pd.ExcelFile(file_path)
            sheet_texts = []
            for sheet in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet)
                sheet_texts.append(f"Sheet: {sheet}")
                sheet_texts.append(df.to_string(index=False))
            text = "\n\n".join(sheet_texts).strip()
            docs = [Document(page_content=text)] if text else []
            return {
                "documents": docs,
                "text": text,
                "ocr_used": False,
                "extraction_mode": "xlsx_text",
                "exception_reason": None if text else "No extractable text in Excel",
            }

    except Exception as e:
        return {
            "documents": [],
            "text": "",
            "ocr_used": False,
            "extraction_mode": "failed",
            "exception_reason": str(e),
        }

    return {
        "documents": [],
        "text": "",
        "ocr_used": False,
        "extraction_mode": "unsupported",
        "exception_reason": f"Unsupported file type: {suffix}",
    }


def create_vectorstore(docs):
    if not docs:
        return None

    splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=150)
    chunks = splitter.split_documents(docs)
    if not chunks:
        return None

    for chunk in chunks:
        chunk.metadata = {"source": st.session_state.get("current_file", "unknown")}

    try:
        emb = get_embeddings(st.session_state["api_key"])
        return Chroma.from_documents(chunks, embedding=emb)
    except Exception:
        return None


def push_agent_log(message):
    st.session_state.agent_logs.append(message)
    refresh_live_batch_activity()


def record_agent_event(step, status, message=""):
    now = time.time()

    if "agent_timings" not in st.session_state:
        st.session_state["agent_timings"] = {}

    if status == "running":
        if step not in st.session_state["agent_timings"]:
            st.session_state["agent_timings"][step] = {}
        if not st.session_state["agent_timings"][step].get("started_at"):
            st.session_state["agent_timings"][step]["started_at"] = now
        st.session_state["active_agent"] = step

    elif status in ["done", "error"]:
        if step not in st.session_state["agent_timings"]:
            st.session_state["agent_timings"][step] = {}
        started_at = st.session_state["agent_timings"][step].get("started_at")
        st.session_state["agent_timings"][step]["ended_at"] = now
        if started_at:
            st.session_state["agent_timings"][step]["elapsed"] = round(now - started_at, 2)
        if st.session_state.get("active_agent") == step:
            st.session_state["active_agent"] = None

    st.session_state.agent_events.append({
        "step": step,
        "status": status,
        "message": message,
    })
    refresh_live_batch_activity()
    
def refresh_live_batch_activity():
    step_placeholder = st.session_state.get("live_step_placeholder")
    progress_placeholder = st.session_state.get("live_progress_placeholder")
    event_placeholder = st.session_state.get("live_event_placeholder")

    total_files = st.session_state.get("batch_total_files", 0)
    processed_files = st.session_state.get("batch_processed_files", 0)
    current_file = st.session_state.get("batch_current_file")
    current_step = st.session_state.get("current_step", "Waiting")
    file_statuses = st.session_state.get("batch_file_statuses", [])
    exception_count = len(st.session_state.get("exception_queue", []))
    per_file_progress = int(st.session_state.get("progress_value", 0))

    if total_files > 0:
        overall_progress = int(((processed_files + (per_file_progress / 100.0)) / total_files) * 100)
        overall_progress = max(0, min(100, overall_progress))
    else:
        overall_progress = per_file_progress

    if step_placeholder is not None:
        elapsed = st.session_state.get("batch_elapsed_seconds", 0.0)

        if total_files > 0:
            elapsed_line = f"**Elapsed:** {elapsed:.2f} sec  " if elapsed > 0 else ""

            step_placeholder.markdown(
                f"""
#### Batch Progress

**Current File:** {current_file or "-"}  
**Current Step:** {current_step}  
**Processed:** {processed_files} / {total_files}  
**Exceptions:** {exception_count}  
{elapsed_line}
"""
            )
        else:
            if current_step != "Waiting":
                step_placeholder.markdown(f"#### Progress\n\n**Current Step:** {current_step}")
            else:
                step_placeholder.empty()

    if progress_placeholder is not None:
        if total_files > 0 or per_file_progress > 0:
            progress_placeholder.progress(overall_progress)
        else:
            progress_placeholder.empty()

    if event_placeholder is not None:
        content = []

        if total_files > 0:
            content.append("#### File Queue")

            if file_statuses:
                for item in file_statuses:
                    status = item.get("status", "pending")
                    file_name = item.get("file_name", "")

                    if status == "done":
                        icon = "✅"
                    elif status == "error":
                        icon = "❌"
                    elif status == "running":
                        icon = "🔄"
                    else:
                        icon = "⏳"

                    line = f"{icon} **{file_name}**"
                    if item.get("message"):
                        line += f"  \n{item.get('message')}"
                    content.append(line)
            else:
                content.append("_No files started yet_")
        else:
            events = st.session_state.get("agent_events", [])
            if events:
                content.append("#### Completed Steps")
                for event in events[-8:]:
                    status = event.get("status", "pending")
                    if status == "done":
                        icon = "✅"
                    elif status == "error":
                        icon = "❌"
                    elif status == "running":
                        icon = "🔄"
                    else:
                        icon = "⏳"

                    line = f"{icon} **{event.get('step', '')}**"
                    if event.get("message"):
                        line += f"  \n{event.get('message')}"
                    content.append(line)

        event_placeholder.markdown("\n\n".join(content) if content else "")

    render_agent_pipeline()



def render_agent_pipeline():
    pipeline_placeholder = st.session_state.get("live_pipeline_placeholder")
    if pipeline_placeholder is None:
        return

    doc_type = st.session_state.get("doc_type")
    events = st.session_state.get("agent_events", [])
    timings = st.session_state.get("agent_timings", {})
    active_agent = st.session_state.get("active_agent")

    pipeline = [
        "Ingestion Agent",
        "Extraction Agent",
        "Retrieval Agent",
        "Classification Agent",
        "Structuring Agent",
        "Validation Agent",
        "Output Agent",
    ]

    if doc_type in ["invoice", "ticket"]:
        pipeline.append("Concur Agent")

    status_map = {name: {"status": "pending", "message": ""} for name in pipeline}

    for event in events:
        step = event.get("step")
        if step in status_map:
            status_map[step] = {
                "status": event.get("status", "pending"),
                "message": event.get("message", ""),
            }

    html_parts = [
        textwrap.dedent("""
        <div style="margin-top:10px;">
            <div style="font-weight:700;font-size:16px;margin-bottom:10px;">
                Agentic Pipeline
            </div>
            <div style="display:flex;flex-wrap:wrap;gap:10px;">
        """).strip()
    ]

    for step in pipeline:
        item = status_map[step]
        status = item["status"]
        elapsed = timings.get(step, {}).get("elapsed")
        running_since = timings.get(step, {}).get("started_at")

        short_name = step.replace(" Agent", "")

        if status == "done":
            bg = "#e8f7ee"
            border = "#8fd19e"
            icon = "✅"
            text = "#166534"
        elif status == "running" or step == active_agent:
            bg = "#eef4ff"
            border = "#7aa2ff"
            icon = "🔄"
            text = "#1d4ed8"
        elif status == "error":
            bg = "#fdecec"
            border = "#f5a3a3"
            icon = "❌"
            text = "#b42318"
        else:
            bg = "#f5f5f5"
            border = "#dddddd"
            icon = "⏳"
            text = "#6b7280"

        if elapsed is not None:
            subtitle = f"{elapsed:.2f}s"
        elif (status == "running" or step == active_agent) and running_since:
            subtitle = f"{round(time.time() - running_since, 2)}s"
        else:
            subtitle = item.get("message") or "Pending"

        card_html = textwrap.dedent(f"""
        <div style="
            min-width:120px;
            flex:1 1 120px;
            padding:12px 10px;
            border-radius:14px;
            border:1px solid {border};
            background:{bg};
            text-align:center;
        ">
            <div style="font-size:18px;line-height:1;">{icon}</div>
            <div style="font-weight:700;color:{text};font-size:12px;margin-top:6px;">
                {short_name}
            </div>
            <div style="font-size:11px;color:#4b5563;margin-top:4px;">
                {subtitle}
            </div>
        </div>
        """).strip()

        html_parts.append(card_html)

    html_parts.append("</div></div>")

    pipeline_placeholder.markdown("".join(html_parts), unsafe_allow_html=True)


def update_batch_file_status(file_name, status, message=""):
    statuses = st.session_state.get("batch_file_statuses", [])

    found = False
    for item in statuses:
        if item.get("file_name") == file_name:
            item["status"] = status
            item["message"] = message
            found = True
            break

    if not found:
        statuses.append({
            "file_name": file_name,
            "status": status,
            "message": message
        })

    st.session_state["batch_file_statuses"] = statuses
    refresh_live_batch_activity()


def update_progress(percent, message):
    st.session_state["progress_value"] = percent
    st.session_state["current_step"] = message

    current_file = st.session_state.get("batch_current_file")
    if current_file:
        update_batch_file_status(current_file, "running", message)

    refresh_live_batch_activity()


def get_suggested_questions(doc_type):
    return [
        "Summarize this architecture document",
        "What systems and components are involved?",
        "What is the design flow?",
        "What assumptions and risks are identified?",
        "What glossary can be extracted from this document?",
    ]


def normalize_graph_result(result):
    if not isinstance(result, dict):
        return {
            "doc_type": None,
            "structured_data": None,
            "result": {},
            "error": "Graph returned non-dict output",
        }

    doc_type = result.get("doc_type") or result.get("type")
    structured_data = result.get("data") if doc_type in ["invoice", "ticket"] else None
    inner = result.get("result", {}) if isinstance(result.get("result", {}), dict) else {}

    return {
        "doc_type": doc_type,
        "structured_data": structured_data,
        "result": inner,
        "error": result.get("error"),
        "step_metrics": result.get("step_metrics", []),
        "confidence": result.get("confidence"),
        "validation": result.get("validation"),
        "ocr_used": result.get("ocr_used", False),
        "extraction_mode": result.get("extraction_mode"),
        "exception_reason": result.get("exception_reason"),
        "needs_review": result.get("needs_review", False),
    }

def process_single_file(uploaded_file):
    reset_single_file_state()
    st.session_state.current_file = uploaded_file.name

    record_agent_event("Ingestion Agent", "running", "Receiving technical document")
    update_progress(5, "Ingestion Agent — file received")
    record_agent_event("Ingestion Agent", "done", "File received")

    record_agent_event("Extraction Agent", "running", "Extracting text")
    extracted = process_file_with_fallback(uploaded_file)
    docs = extracted["documents"]
    full_text = extracted["text"]

    if extracted["ocr_used"]:
        record_agent_event("Extraction Agent", "done", "Text extracted using OCR fallback")
    else:
        record_agent_event("Extraction Agent", "done", "Text extracted")

    update_progress(20, "Extraction Agent — text extracted")

    if not full_text:
        reason = extracted["exception_reason"] or "No extractable text"
        return {
            "file_name": uploaded_file.name,
            "status": "Exception",
            "doc_type": "technical_doc",
            "ocr_used": extracted["ocr_used"],
            "exception_reason": reason,
            "review_data": None,
            "validation": None,
            "confidence": None,
            "auto_result": None,
            "vectorstore": None,
            "full_text": None,
            "cost": 0.0,
            "tokens": 0,
        }

    st.session_state.full_text = full_text

    record_agent_event("Retrieval Agent", "running", "Creating vector index")
    vectorstore = create_vectorstore(docs)
    st.session_state.vectorstore = vectorstore
    record_agent_event("Retrieval Agent", "done", "Vector index created")
    update_progress(30, "Retrieval Agent — search index ready")

    graph = build_graph()
    graph_input = {
        "text": full_text,
        "filename": uploaded_file.name,
        "template": get_active_template_bytes(),
        "progress": update_progress,
        "event_callback": record_agent_event,
        "ocr_used": extracted["ocr_used"],
        "extraction_mode": extracted["extraction_mode"],
        "exception_reason": extracted["exception_reason"],
    }

    before_cost = st.session_state["metrics"]["cost"]
    before_tokens = st.session_state["metrics"]["tokens"]

    raw_result = graph.invoke(graph_input)
    normalized = normalize_graph_result(raw_result)

    doc_type = "technical_doc"
    result = normalized.get("result", {})
    review_data = result.get("data") or normalized.get("structured_data") or {}

    record_agent_event("Validation Agent", "running", "Checking extracted technical content")
    validation = validate_document_data(review_data, doc_type)
    confidence = build_confidence_map(review_data, doc_type)
    record_agent_event("Validation Agent", "done", "Best-effort validation completed")

    exception_reason = classify_exception(
        doc_type=doc_type,
        text=full_text,
        validation=validation,
        confidence=confidence,
        extraction_meta=extracted,
    )

    st.session_state.doc_type = doc_type
    st.session_state.review_data = review_data
    st.session_state.validation_result = validation
    st.session_state.confidence_map = confidence
    st.session_state.auto_result = {
        "doc_type": doc_type,
        "structured_data": normalized.get("structured_data"),
        "result": result,
        "metrics": {},
        "step_metrics": normalized.get("step_metrics", []),
        "ocr_used": extracted["ocr_used"],
        "extraction_mode": extracted["extraction_mode"],
    }
    st.session_state.suggested_questions = get_suggested_questions(doc_type)

    after_cost = st.session_state["metrics"]["cost"]
    after_tokens = st.session_state["metrics"]["tokens"]

    status = "Completed"
    if extracted.get("exception_reason") and not full_text:
        status = "Exception"

    record_agent_event("Workflow Agent", "done", "Technical document processing completed")
    update_progress(100, "Workflow Agent — completed")

    return {
        "file_name": uploaded_file.name,
        "status": status,
        "doc_type": doc_type,
        "ocr_used": extracted["ocr_used"],
        "exception_reason": exception_reason,
        "review_data": review_data,
        "validation": validation,
        "confidence": confidence,
        "auto_result": st.session_state.auto_result,
        "vectorstore": vectorstore,
        "full_text": full_text,
        "cost": round(after_cost - before_cost, 6),
        "tokens": after_tokens - before_tokens,
    }

def load_batch_result_into_session(index):
    if index < 0 or index >= len(st.session_state.batch_results):
        return

    item = st.session_state.batch_results[index]
    st.session_state.active_batch_index = index
    st.session_state.current_file = item.get("file_name")
    st.session_state.doc_type = item.get("doc_type")
    st.session_state.review_data = item.get("review_data")
    st.session_state.validation_result = item.get("validation")
    st.session_state.confidence_map = item.get("confidence")
    st.session_state.auto_result = item.get("auto_result")
    st.session_state.vectorstore = item.get("vectorstore")
    st.session_state.full_text = item.get("full_text")


def get_batch_signature(uploaded_files):
    if not uploaded_files:
        return None

    parts = []
    for file in uploaded_files:
        try:
            content_hash = hashlib.md5(file.getvalue()).hexdigest()
        except Exception:
            content_hash = f"{file.name}-{len(file.getvalue())}"
        parts.append(f"{file.name}:{content_hash}")

    return "|".join(parts)


def go_to_next_batch_result():
    batch_results = st.session_state.get("batch_results", [])
    if not batch_results:
        return

    current_index = st.session_state.get("active_batch_index", 0)
    next_index = current_index + 1

    if next_index < len(batch_results):
        load_batch_result_into_session(next_index)


def compact_field(label, value):
    st.markdown(
        f"**{label}**  \n<small>{value if value not in [None, ''] else '-'}</small>",
        unsafe_allow_html=True
    )

# ------------------------------
# REVIEW / ACTIONS
# ------------------------------
def render_validation_summary():
    validation = st.session_state.get("validation_result") or {}
    warnings = validation.get("warnings", [])

    st.markdown("#### Extraction Notes")
    st.success("Architecture report generated using best available extraction")

    for item in warnings:
        st.caption(f"• {item}")


def render_confidence_table():
    confidence = st.session_state.get("confidence_map") or {}
    if not confidence:
        return

    rows = [{"Field": k, "Confidence": v.get("label", "-"), "Reason": v.get("reason", "-")} for k, v in confidence.items()]
    st.markdown("#### Confidence")
    st.dataframe(pd.DataFrame(rows), use_container_width=True, height=220, hide_index=True)


# ------------------------------
# UI
# ------------------------------
def render_header():
    logo_path = Path(__file__).parent / "TDReader.png"
    col_logo, col_title = st.columns([1, 6], gap="small")

    with col_logo:
        if logo_path.exists():
            st.image(logo_path, width=500)

    with col_title:
        st.markdown("## Technical Document Reader")
        st.caption("AI-powered architecture, design, and specification understanding")


def render_sidebar_and_upload():
    with st.sidebar:
        st.write(f"Hi **{st.session_state['user']}**")

        st.markdown("---")

        model_choice = st.selectbox(
            "Choose Model",
            ["gpt-4o-mini", "gpt-4o", "gpt-5"],
            index=["gpt-4o-mini", "gpt-4o", "gpt-5"].index(
                st.session_state.get("model_choice", "gpt-4o-mini")
            )
        )
        st.session_state["model_choice"] = model_choice

        st.markdown("---")
        st.success("🔑 API key loaded securely")
        cost = st.session_state.get("metrics", {}).get("cost", 0.0)

        st.markdown("---")
        st.write(f"💰 Session Cost ${round(cost, 6)}")

        st.markdown("---")
        if st.button("Logout", use_container_width=True):
            for key in ["logged_in", "user", "role", "api_key"]:
                if key in st.session_state:
                    del st.session_state[key]
            st.rerun()

    c1, c2 = st.columns([6, 1], gap="small")
    with c1:
        uploaded_files = st.file_uploader(
            f"Upload technical documents - max {MAX_BATCH_FILES} files per batch",
            type=["txt", "pdf", "docx", "pptx", "png", "jpg", "jpeg"],
            accept_multiple_files=True,
            key=f"main_file_uploader_{st.session_state.uploader_key}"
        )

    with c2:
        st.write("")
        st.write("")
        if st.button("Reset", use_container_width=True):
            st.session_state.batch_results = []
            st.session_state.exception_queue = []
            st.session_state.batch_processed = False
            st.session_state.last_batch_signature = None
            st.session_state.show_reprocess_confirm = False
            st.session_state.pending_batch_signature = None
            st.session_state.batch_total_files = 0
            st.session_state.batch_processed_files = 0
            st.session_state.batch_current_file = None
            st.session_state.batch_file_statuses = []
            st.session_state.uploader_key += 1
            reset_run_state()
            st.rerun()

    if uploaded_files and len(uploaded_files) > MAX_BATCH_FILES:
        st.error(f"Batch limit exceeded. Maximum allowed is {MAX_BATCH_FILES} files.")
        uploaded_files = uploaded_files[:MAX_BATCH_FILES]

    st.markdown("---")
    return uploaded_files


def render_result_workspace():
    st.markdown("### Result")

    if not st.session_state.get("auto_result"):
        st.caption("Process a technical document to view results.")
        return

    doc_type = st.session_state.get("doc_type")
    result = st.session_state.get("auto_result", {}).get("result", {})
    data = st.session_state.get("review_data") or {}

    current_index = st.session_state.get("active_batch_index", 0)
    total_results = len(st.session_state.get("batch_results", []))
    has_next = current_index < (total_results - 1)

    if doc_type != "technical_doc":
        st.caption("No technical document result available.")
        return

    st.markdown("#### Professional Overview")
    st.write(
        data.get("executive_overview")
        or "Overview could not be fully extracted. Best available technical details are shown below."
    )

    st.markdown("#### Architecture Profile")
    st.caption(f"Architecture Style: {data.get('architecture_style') or '-'}")
    st.caption(f"Deployment Model: {data.get('deployment_model') or '-'}")
    platforms = data.get("primary_platforms", [])
    st.caption(f"Primary Platforms: {', '.join(platforms) if platforms else '-'}")

    st.markdown("#### Visual Flow")
    st.code(data.get("visual_flow") or "No visual flow extracted", language="text")

    c1, c2 = st.columns(2)

    with c1:
        st.markdown("#### Systems")
        systems = data.get("systems", [])
        if systems:
            for item in systems:
                st.caption(f"• {item}")
        else:
            st.caption("• No systems identified")

        st.markdown("#### Components")
        components = data.get("components", [])
        if components:
            for item in components:
                st.caption(f"• {item}")
        else:
            st.caption("• No components identified")

        st.markdown("#### Interfaces")
        interfaces = data.get("interfaces", [])
        if interfaces:
            for item in interfaces:
                st.caption(f"• {item}")
        else:
            st.caption("• No interfaces identified")

        st.markdown("#### Actors")
        actors = data.get("actors", [])
        if actors:
            for item in actors:
                st.caption(f"• {item}")
        else:
            st.caption("• No actors identified")

        st.markdown("#### Data Entities")
        entities = data.get("data_entities", [])
        if entities:
            for item in entities:
                st.caption(f"• {item}")
        else:
            st.caption("• No data entities identified")

    with c2:
        st.markdown("#### Design Flow")
        flow = data.get("design_flow", [])
        if flow:
            for idx, item in enumerate(flow, start=1):
                st.caption(f"{idx}. {item}")
        else:
            st.caption("• No design flow extracted")

        st.markdown("#### Dependencies")
        deps = data.get("dependencies", [])
        if deps:
            for item in deps:
                st.caption(f"• {item}")
        else:
            st.caption("• No dependencies identified")

        st.markdown("#### Integration Points")
        ints = data.get("integration_points", [])
        if ints:
            for item in ints:
                st.caption(f"• {item}")
        else:
            st.caption("• No integration points identified")

        st.markdown("#### Monitoring and Observability")
        obs = data.get("monitoring_observability", [])
        if obs:
            for item in obs:
                st.caption(f"• {item}")
        else:
            st.caption("• No monitoring details identified")

    render_validation_summary()
    render_confidence_table()

    with st.expander("Glossary", expanded=False):
        glossary = data.get("glossary", [])
        if glossary:
            rows = [{"Term": g.get("term", ""), "Meaning": g.get("meaning", "")} for g in glossary]
            st.dataframe(pd.DataFrame(rows), use_container_width=True, hide_index=True)
        else:
            st.caption("No glossary items extracted")

    with st.expander("Requirements, Assumptions, Risks", expanded=False):
        st.markdown("**Functional Requirements**")
        for item in data.get("functional_requirements", []):
            st.caption(f"• {item}")
        if not data.get("functional_requirements"):
            st.caption("• None extracted")

        st.markdown("**Non-Functional Requirements**")
        for item in data.get("non_functional_requirements", []):
            st.caption(f"• {item}")
        if not data.get("non_functional_requirements"):
            st.caption("• None extracted")

        st.markdown("**Security Considerations**")
        for item in data.get("security_considerations", []):
            st.caption(f"• {item}")
        if not data.get("security_considerations"):
            st.caption("• None extracted")

        st.markdown("**Assumptions**")
        for item in data.get("assumptions", []):
            st.caption(f"• {item}")
        if not data.get("assumptions"):
            st.caption("• None extracted")

        st.markdown("**Constraints**")
        for item in data.get("constraints", []):
            st.caption(f"• {item}")
        if not data.get("constraints"):
            st.caption("• None extracted")

        st.markdown("**Risks**")
        for item in data.get("risks", []):
            st.caption(f"• {item}")
        if not data.get("risks"):
            st.caption("• None extracted")

        st.markdown("**Recommendations**")
        for item in data.get("recommendations", []):
            st.caption(f"• {item}")
        if not data.get("recommendations"):
            st.caption("• None extracted")

        st.markdown("**Open Questions**")
        for item in data.get("open_questions", []):
            st.caption(f"• {item}")
        if not data.get("open_questions"):
            st.caption("• None extracted")

    t1, t2, t3, t4 = st.columns(4)

    with t1:
        summary_md = result.get("summary_markdown")
        if summary_md:
            st.download_button(
                "Download MD",
                data=summary_md.encode("utf-8"),
                file_name=result.get("file_name", "technical_architecture_report.md"),
                mime="text/markdown",
                use_container_width=True
            )

    with t2:
        summary_pdf = result.get("summary_pdf")
        if summary_pdf:
            st.download_button(
                "Download PDF",
                data=summary_pdf,
                file_name=result.get("pdf_file_name", "technical_architecture_report.pdf"),
                mime="application/pdf",
                use_container_width=True
            )

    with t3:
        summary_docx = result.get("summary_docx")
        if summary_docx:
            st.download_button(
                "Download DOCX",
                data=summary_docx,
                file_name=result.get("docx_file_name", "technical_architecture_report.docx"),
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True
            )

    with t4:
        if st.button("Next Document", use_container_width=True, disabled=not has_next, key="technical_doc_next"):
            go_to_next_batch_result()
            st.rerun()


def render_batch_table():
    st.markdown("### Batch Results")
    elapsed = st.session_state.get("batch_elapsed_seconds", 0.0)
    if elapsed:
        st.caption(f"Batch processed in {elapsed:.2f} sec")

    if not st.session_state.batch_results:
        st.caption("No batch results yet")
        return

    rows = []
    for item in st.session_state.batch_results:
        rows.append({
            "File": item.get("file_name"),
            "Type": item.get("doc_type"),
            "Status": item.get("status"),
            "OCR": "Yes" if item.get("ocr_used") else "No",
            "Cost": item.get("cost"),
            "Tokens": item.get("tokens"),
        })

    df = pd.DataFrame(rows)
    st.dataframe(df, use_container_width=True, hide_index=True, height=220)

    current_index = st.session_state.get("active_batch_index", 0)
    if st.session_state.batch_results:
        current_index = max(0, min(current_index, len(st.session_state.batch_results) - 1))

    selected = st.selectbox(
        "Open processed document",
        options=list(range(len(st.session_state.batch_results))),
        format_func=lambda i: f"{st.session_state.batch_results[i]['file_name']} ({st.session_state.batch_results[i]['status']})",
        index=current_index,
    )
    if selected is not None:
        load_batch_result_into_session(selected)


def render_exception_queue():
    st.markdown("### Processing Notes")
    if not st.session_state.exception_queue:
        st.caption("No processing exceptions")
        return

    rows = []
    for item in st.session_state.exception_queue:
        rows.append({
            "File": item.get("file_name"),
            "Type": item.get("doc_type"),
            "Reason": item.get("exception_reason"),
            "OCR": "Yes" if item.get("ocr_used") else "No",
        })

    df = pd.DataFrame(rows)
    st.dataframe(df, use_container_width=True, hide_index=True, height=200)


def render_template_manager():
    st.markdown("### Template Manager")

    template_upload = st.file_uploader(
        "Upload Architecture / HLD / LLD / SRS Template",
        type=["docx"],
        key="template_manager_uploader"
    )

    if template_upload and st.button("Add Template", use_container_width=True):
        add_template_to_library(template_upload)
        st.success("Template added to library")
        st.rerun()

    library = st.session_state.get("template_library", [])
    if not library:
        st.caption("No template uploaded. The app will generate the report in the best possible format.")
        return

    selected = st.selectbox(
        "Choose active template",
        options=list(range(len(library))),
        format_func=lambda i: library[i]["name"],
        index=st.session_state.get("active_template_index", 0) if library else 0,
        key="active_template_selector"
    )
    st.session_state.active_template_index = selected

    active = library[selected]
    st.success(f"Active template: {active['name']}")

def render_batch_downloads():
    st.markdown("### Batch Downloads")
    st.caption("Batch download is not enabled for this technical-document version.")

# ------------------------------
# MAIN
# ------------------------------
render_header()
uploaded_files = render_sidebar_and_upload()

left_col, right_col = st.columns([1, 1.6], gap="large")

with left_col:
    st.markdown("### Activity")
    st.session_state["live_step_placeholder"] = st.empty()
    st.session_state["live_progress_placeholder"] = st.empty()
    st.session_state["live_event_placeholder"] = st.empty()
    refresh_live_batch_activity()

    st.markdown("---")
    render_agent_pipeline()

    current_batch_signature = get_batch_signature(uploaded_files)
    last_batch_signature = st.session_state.get("last_batch_signature")

    process_disabled = not uploaded_files

    if st.button("Process Batch", use_container_width=True, disabled=process_disabled):
        if current_batch_signature and current_batch_signature == last_batch_signature:
            st.session_state.show_reprocess_confirm = True
            st.session_state.pending_batch_signature = current_batch_signature
        else:
            st.session_state.batch_results = []
            st.session_state.exception_queue = []
            st.session_state.show_reprocess_confirm = False
            st.session_state.pending_batch_signature = None

            st.session_state.batch_started_at = time.time()
            st.session_state.batch_completed_at = None
            st.session_state.batch_elapsed_seconds = 0.0

            st.session_state.batch_total_files = len(uploaded_files)
            st.session_state.batch_processed_files = 0
            st.session_state.batch_current_file = None
            st.session_state.batch_file_statuses = [
                {"file_name": f.name, "status": "pending", "message": ""}
                for f in uploaded_files
            ]
            refresh_live_batch_activity()

            for uploaded_file in uploaded_files:
                st.session_state.batch_current_file = uploaded_file.name
                update_batch_file_status(uploaded_file.name, "running", "Processing started")
                refresh_live_batch_activity()

                result = process_single_file(uploaded_file)
                st.session_state.batch_results.append(result)

                if result.get("status") == "Exception":
                    st.session_state.exception_queue.append(result)
                    update_batch_file_status(
                        uploaded_file.name,
                        "error",
                        result.get("exception_reason", "Exception")
                    )
                else:
                    update_batch_file_status(uploaded_file.name, "done", result.get("status", "Completed"))

                st.session_state.batch_processed_files += 1
                st.session_state["progress_value"] = 0
                refresh_live_batch_activity()

            if st.session_state.batch_results:
                load_batch_result_into_session(0)
                st.session_state.batch_processed = True
                st.session_state.last_batch_signature = current_batch_signature
                st.session_state.batch_completed_at = time.time()
                st.session_state.batch_elapsed_seconds = (
                    st.session_state.batch_completed_at - st.session_state.batch_started_at
                )
                st.success("Batch processing completed")

with right_col:
    render_result_workspace()

st.markdown("---")
render_batch_table()
render_exception_queue()
render_batch_downloads()

st.markdown("---")
render_template_manager()

with st.expander("Metrics", expanded=False):
    m = st.session_state.get("metrics", {})
    c1, c2, c3, c4 = st.columns(4)
    c1.metric("Total Cost", f"${m.get('cost', 0.0):.6f}")
    c2.metric("Total Tokens", m.get("tokens", 0))
    c3.metric("Input Tokens", m.get("input_tokens", 0))
    c4.metric("Output Tokens", m.get("output_tokens", 0))

    doc_costs = st.session_state.get("doc_costs", {})
    doc_rows = [
        {"Document": k, "Cost": round(v.get("cost", 0.0), 6), "Tokens": v.get("tokens", 0)}
        for k, v in doc_costs.items()
    ]
    if doc_rows:
        st.dataframe(pd.DataFrame(doc_rows), use_container_width=True, hide_index=True, height=220)
